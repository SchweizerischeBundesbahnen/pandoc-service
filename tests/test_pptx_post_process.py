import io
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

from app import PptxPostProcess
from app.PptxPostProcess import SLIDE_SIZES, _apply_slide_size, process


def create_test_pptx() -> bytes:
    """Create a minimal test PPTX file."""
    prs = Presentation()
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_slide_layout)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def test_process_with_no_slide_size():
    """Test processing PPTX without specifying slide size (no modifications)."""
    pptx_bytes = create_test_pptx()
    original_prs = Presentation(io.BytesIO(pptx_bytes))
    original_width = original_prs.slide_width
    original_height = original_prs.slide_height

    result_bytes = process(pptx_bytes, slide_size=None)
    result_prs = Presentation(io.BytesIO(result_bytes))

    # Dimensions should remain unchanged
    assert result_prs.slide_width == original_width
    assert result_prs.slide_height == original_height


def test_process_with_16_9_slide_size():
    """Test processing PPTX with 16:9 slide size."""
    pptx_bytes = create_test_pptx()

    result_bytes = process(pptx_bytes, slide_size="16:9")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["16:9"]["width"])
    expected_height = Inches(SLIDE_SIZES["16:9"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_process_with_4_3_slide_size():
    """Test processing PPTX with 4:3 slide size."""
    pptx_bytes = create_test_pptx()

    result_bytes = process(pptx_bytes, slide_size="4:3")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["4:3"]["width"])
    expected_height = Inches(SLIDE_SIZES["4:3"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_process_with_case_insensitive_slide_size():
    """Test that slide size is case insensitive."""
    pptx_bytes = create_test_pptx()

    # Test lowercase
    result_bytes = process(pptx_bytes, slide_size="a4")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["A4"]["width"])
    expected_height = Inches(SLIDE_SIZES["A4"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_process_with_invalid_slide_size():
    """Test processing PPTX with invalid slide size raises ValueError."""
    pptx_bytes = create_test_pptx()

    with pytest.raises(ValueError, match="Unsupported slide size: INVALID"):
        process(pptx_bytes, slide_size="INVALID")


def test_process_with_letter_slide_size():
    """Test processing PPTX with LETTER slide size."""
    pptx_bytes = create_test_pptx()

    result_bytes = process(pptx_bytes, slide_size="LETTER")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["LETTER"]["width"])
    expected_height = Inches(SLIDE_SIZES["LETTER"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_process_with_a3_slide_size():
    """Test processing PPTX with A3 slide size."""
    pptx_bytes = create_test_pptx()

    result_bytes = process(pptx_bytes, slide_size="A3")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["A3"]["width"])
    expected_height = Inches(SLIDE_SIZES["A3"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_process_with_widescreen_slide_size():
    """Test processing PPTX with WIDESCREEN slide size."""
    pptx_bytes = create_test_pptx()

    result_bytes = process(pptx_bytes, slide_size="WIDESCREEN")
    result_prs = Presentation(io.BytesIO(result_bytes))

    expected_width = Inches(SLIDE_SIZES["WIDESCREEN"]["width"])
    expected_height = Inches(SLIDE_SIZES["WIDESCREEN"]["height"])

    assert result_prs.slide_width == expected_width
    assert result_prs.slide_height == expected_height


def test_apply_slide_size_none():
    """Test _apply_slide_size with None does nothing."""
    prs = Presentation()
    original_width = prs.slide_width
    original_height = prs.slide_height

    _apply_slide_size(prs, slide_size=None)

    # Dimensions should remain unchanged
    assert prs.slide_width == original_width
    assert prs.slide_height == original_height


def test_apply_slide_size_with_valid_size():
    """Test _apply_slide_size with valid slide size."""
    prs = Presentation()

    _apply_slide_size(prs, slide_size="16:9")

    expected_width = Inches(SLIDE_SIZES["16:9"]["width"])
    expected_height = Inches(SLIDE_SIZES["16:9"]["height"])

    assert prs.slide_width == expected_width
    assert prs.slide_height == expected_height


def test_apply_slide_size_with_invalid_size():
    """Test _apply_slide_size with invalid slide size raises ValueError."""
    prs = Presentation()

    with pytest.raises(ValueError, match="Unsupported slide size: BADSIZE"):
        _apply_slide_size(prs, slide_size="BADSIZE")


def test_slide_sizes_constant():
    """Test that SLIDE_SIZES constant has expected entries."""
    expected_sizes = ["16:9", "WIDESCREEN", "4:3", "LETTER", "LEDGER", "A4", "A3"]

    for size in expected_sizes:
        assert size in SLIDE_SIZES
        assert "width" in SLIDE_SIZES[size]
        assert "height" in SLIDE_SIZES[size]
        assert isinstance(SLIDE_SIZES[size]["width"], (int, float))
        assert isinstance(SLIDE_SIZES[size]["height"], (int, float))


def test_process_returns_bytes():
    """Test that process returns bytes."""
    pptx_bytes = create_test_pptx()

    result = process(pptx_bytes, slide_size="16:9")

    assert isinstance(result, bytes)
    assert len(result) > 0


def test_process_preserves_content():
    """Test that processing preserves presentation content."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add a title shape to the slide
    left = top = Inches(1)
    width = height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.text = "Test Content"

    out = io.BytesIO()
    prs.save(out)
    pptx_bytes = out.getvalue()

    # Process the PPTX
    result_bytes = process(pptx_bytes, slide_size="16:9")
    result_prs = Presentation(io.BytesIO(result_bytes))

    # Check that content is preserved
    assert len(result_prs.slides) == 1
    result_slide = result_prs.slides[0]
    # Find the textbox in the result
    textboxes = [shape for shape in result_slide.shapes if hasattr(shape, "text")]
    assert len(textboxes) > 0
    assert textboxes[0].text == "Test Content"
